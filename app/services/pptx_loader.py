"""
PPTX Loader - Preserves slide structure, formatting, and metadata
"""
from typing import List, Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PptxLoader:
    """
    Loads PPTX files and extracts text while preserving:
    - Slide structure
    - Shape positions and formatting
    - Tables structure
    - Images (metadata)
    - Text formatting
    """

    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Load PPTX and extract structured content.

        Returns:
            List of slide dicts, each containing:
            - slide_index: int
            - shapes: List of shape data with text and metadata
        """
        prs = Presentation(file_path)
        slides_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                'slide_index': slide_idx,
                'shapes': []
            }

            for shape_idx, shape in enumerate(slide.shapes):
                shape_data = self._extract_shape_data(shape, shape_idx)
                if shape_data:
                    slide_data['shapes'].append(shape_data)

            slides_data.append(slide_data)

        return slides_data

    def _extract_shape_data(self, shape, shape_idx: int) -> Dict[str, Any] | None:
        """Extract data from a single shape."""
        shape_data = {
            'shape_index': shape_idx,
            'shape_type': shape.shape_type,
            'name': shape.name if hasattr(shape, 'name') else '',
        }

        # Extract text from text-containing shapes
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            shape_data['type'] = 'text'
            shape_data['paragraphs'] = self._extract_paragraphs(shape.text_frame)
            return shape_data

        # Extract table content
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_data['type'] = 'table'
            shape_data['table_data'] = self._extract_table(shape.table)
            return shape_data

        # Track images (we'll preserve them without translation)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_data['type'] = 'image'
            return shape_data

        # Track groups (contain multiple shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_data['type'] = 'group'
            shape_data['shapes'] = []
            for sub_idx, sub_shape in enumerate(shape.shapes):
                sub_data = self._extract_shape_data(sub_shape, sub_idx)
                if sub_data:
                    shape_data['shapes'].append(sub_data)
            return shape_data

        return None

    def _extract_paragraphs(self, text_frame) -> List[Dict[str, Any]]:
        """Extract paragraphs with formatting metadata."""
        paragraphs_data = []

        for para in text_frame.paragraphs:
            para_data = {
                'text': para.text,
                'level': para.level,
                'alignment': para.alignment,
                'runs': []
            }

            # Extract run-level formatting
            for run in para.runs:
                run_data = {
                    'text': run.text,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                    'font_size': run.font.size.pt if run.font.size else None,
                    'font_name': run.font.name,
                }
                para_data['runs'].append(run_data)

            paragraphs_data.append(para_data)

        return paragraphs_data

    def _extract_table(self, table) -> Dict[str, Any]:
        """Extract table content with structure."""
        rows_data = []

        for row in table.rows:
            cells_data = []
            for cell in row.cells:
                cell_text = cell.text
                cells_data.append({
                    'text': cell_text,
                    'has_text_frame': hasattr(cell, 'text_frame')
                })
            rows_data.append(cells_data)

        return {
            'rows': rows_data,
            'num_rows': len(table.rows),
            'num_cols': len(table.columns)
        }
