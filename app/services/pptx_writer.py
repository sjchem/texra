"""
PPTX Writer - Preserves slide structure, formatting, and metadata
"""
import tempfile
from typing import List, Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


class PptxWriter:
    """
    Writes translated text back to PPTX while preserving:
    - Original slide structure
    - Shape positions and formatting
    - Tables structure
    - Images
    - Text formatting (bold, italic, font size, etc.)
    """

    def write(self, slides_data: List[Dict[str, Any]], original_path: str) -> str:
        """
        Write translated slides back to PPTX format.

        Args:
            slides_data: List of slide dicts with translated text
            original_path: Path to original PPTX (used as template)

        Returns:
            Path to the output PPTX file
        """
        # Load the original presentation as a template
        prs = Presentation(original_path)

        # Update each slide with translated content
        for slide_data in slides_data:
            slide_idx = slide_data['slide_index']

            # Safety check
            if slide_idx >= len(prs.slides):
                continue

            slide = prs.slides[slide_idx]
            self._update_slide(slide, slide_data['shapes'])

        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            output_path = tmp.name

        prs.save(output_path)
        return output_path

    def _update_slide(self, slide, shapes_data: List[Dict[str, Any]]):
        """Update all shapes on a slide with translated content."""
        for shape_data in shapes_data:
            shape_idx = shape_data['shape_index']

            # Safety check
            if shape_idx >= len(slide.shapes):
                continue

            shape = slide.shapes[shape_idx]
            self._update_shape(shape, shape_data)

    def _update_shape(self, shape, shape_data: Dict[str, Any]):
        """Update a single shape with translated content."""
        shape_type = shape_data.get('type')

        if shape_type == 'text':
            self._update_text_shape(shape, shape_data)
        elif shape_type == 'table':
            self._update_table_shape(shape, shape_data)
        elif shape_type == 'group':
            self._update_group_shape(shape, shape_data)
        # Images are preserved as-is

    def _update_text_shape(self, shape, shape_data: Dict[str, Any]):
        """Update text in a shape while preserving formatting."""
        if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        paragraphs_data = shape_data.get('paragraphs', [])

        # Clear existing paragraphs (keep first one)
        for i in range(len(text_frame.paragraphs) - 1, 0, -1):
            p = text_frame.paragraphs[i]._element
            p.getparent().remove(p)

        # Update or create paragraphs
        for para_idx, para_data in enumerate(paragraphs_data):
            if para_idx == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            # Set paragraph properties
            para.text = ""  # Clear existing text
            para.level = para_data.get('level', 0)
            if para_data.get('alignment') is not None:
                para.alignment = para_data['alignment']

            # Add runs with formatting
            runs_data = para_data.get('runs', [])
            if runs_data:
                for run_data in runs_data:
                    run = para.add_run()
                    run.text = run_data.get('text', '')

                    # Apply formatting
                    if run_data.get('bold') is not None:
                        run.font.bold = run_data['bold']
                    if run_data.get('italic') is not None:
                        run.font.italic = run_data['italic']
                    if run_data.get('underline') is not None:
                        run.font.underline = run_data['underline']
                    if run_data.get('font_size') is not None:
                        run.font.size = Pt(run_data['font_size'])
                    if run_data.get('font_name'):
                        run.font.name = run_data['font_name']
            else:
                # Fallback: just use paragraph text
                para.text = para_data.get('text', '')

    def _update_table_shape(self, shape, shape_data: Dict[str, Any]):
        """Update table content while preserving structure."""
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            return

        table = shape.table
        table_data = shape_data.get('table_data', {})
        rows_data = table_data.get('rows', [])

        # Update each cell
        for row_idx, row_data in enumerate(rows_data):
            if row_idx >= len(table.rows):
                break

            for col_idx, cell_data in enumerate(row_data):
                if col_idx >= len(table.columns):
                    break

                cell = table.cell(row_idx, col_idx)
                cell.text = cell_data.get('text', '')

    def _update_group_shape(self, shape, shape_data: Dict[str, Any]):
        """Update shapes within a group."""
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            return

        sub_shapes_data = shape_data.get('shapes', [])

        for sub_shape_data in sub_shapes_data:
            sub_idx = sub_shape_data['shape_index']
            if sub_idx >= len(shape.shapes):
                continue

            sub_shape = shape.shapes[sub_idx]
            self._update_shape(sub_shape, sub_shape_data)
