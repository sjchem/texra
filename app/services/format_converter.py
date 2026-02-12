"""
Format Converter - Convert documents to PDF format
"""
import re
import tempfile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from odf import text as odf_text
from odf.opendocument import load as odf_load
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak, Table, TableStyle
from reportlab.lib import colors


class FormatConverter:
    """Convert various document formats to PDF."""

    @staticmethod
    def docx_to_pdf(docx_path: str) -> str:
        """
        Convert DOCX to PDF.
        Simple conversion that preserves basic text structure.
        """
        doc = Document(docx_path)

        # Create output PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            output_path = tmp.name

        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Convert paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                story.append(Paragraph(para.text, styles['Normal']))

        # Convert tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)

            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                ]))
                story.append(t)

        if not story:
            story.append(Paragraph("(empty document)", styles['Normal']))

        pdf_doc.build(story)
        return output_path

    @staticmethod
    def pptx_to_pdf(pptx_path: str) -> str:
        """
        Convert PPTX to PDF.
        Uses PyMuPDF to create PDF from slide images.
        """
        prs = Presentation(pptx_path)

        # Create output PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            output_path = tmp.name

        pdf_doc = fitz.open()

        # For now, we'll create a simple text-based PDF from slides
        # A full implementation would render slides as images
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet

        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        for slide_idx, slide in enumerate(prs.slides):
            # Add slide number
            story.append(Paragraph(f"<b>Slide {slide_idx + 1}</b>", styles['Heading1']))

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    story.append(Paragraph(shape.text, styles['Normal']))

            if slide_idx < len(prs.slides) - 1:
                story.append(PageBreak())

        if not story:
            story.append(Paragraph("(empty presentation)", styles['Normal']))

        doc.build(story)
        pdf_doc.close()

        return output_path

    @staticmethod
    def odt_to_pdf(odt_path: str) -> str:
        """Convert ODT documents to PDF using a simple text rendering pipeline."""
        odt_doc = odf_load(odt_path)
        paragraphs = odt_doc.getElementsByType(odf_text.P)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            output_path = tmp.name

        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        for para in paragraphs:
            para_str = str(para)
            cleaned = re.sub(r'<[^>]+>', '', para_str).strip()
            if cleaned:
                story.append(Paragraph(cleaned, styles['Normal']))

        if not story:
            story.append(Paragraph("(empty document)", styles['Normal']))

        pdf_doc.build(story)
        return output_path

    @staticmethod
    def to_pdf(file_path: str) -> str:
        """Convert supported formats to PDF, returning the resulting path."""
        lower = file_path.lower()

        if lower.endswith('.docx'):
            return FormatConverter.docx_to_pdf(file_path)
        if lower.endswith('.pptx'):
            return FormatConverter.pptx_to_pdf(file_path)
        if lower.endswith('.odt'):
            return FormatConverter.odt_to_pdf(file_path)
        if lower.endswith('.pdf'):
            return file_path

        raise ValueError(f"Unsupported conversion format for preview: {file_path}")
